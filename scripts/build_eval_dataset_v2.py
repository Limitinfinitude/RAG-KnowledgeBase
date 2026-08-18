# -*- coding: utf-8 -*-
"""构建 v2 评测数据集（权威公开数据为主力 + 自建小文档对照）。

来源与结构：
  1. DuReader-retrieval dev（2000 条带正/负段落标注）→ 抽样查询，正段落拼成多主题 .txt 文档
  2. CMRC 2018（人工标注篇章问答）→ 篇章聚合为长 .md 文档，标注答案作 ground truth
  3. 中文维基长文（0xDing/wikipedia-cn-20230720-filtered）→ 完整长条目 + LLM 生成问答（自检答案在原文）
  4. 自建小文档（eval_corpus_v2/*_20260819.*，含生成的 docx/xlsx/pptx）+ 旧 2026-08-17 十篇
  5. 负样本：纯负（主题不在库，自检裁剪）+ 陷阱（词形相近语义无关）

输出：eval_corpus_v2/ 语料 + scripts/EVAL_SET_V2.json + 语料来源说明。
用法（项目根目录）：conda run -n rag_demo python -X utf8 scripts/build_eval_dataset_v2.py
"""
from __future__ import annotations

import argparse
import gzip
import io
import json
import os
import random
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

RAW = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "eval_raw_data")
OUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "eval_corpus_v2")
OUT_JSON = os.path.join(os.path.dirname(os.path.abspath(__file__)), "EVAL_SET_V2.json")

# 8 字符以上的粗筛禁词（公开语料偶有低质/敏感内容）
_DROP_PATTERNS = re.compile(
    "(性爱|做爱|自慰|阴茎|阴道|性交|赌博|赌场|毒品|枪支|法轮|六四|台独|博彩)"
)


def _clean_text(t: str) -> str:
    t = re.sub(r"\s+", " ", str(t or "")).strip()
    return t


def _passage_ok(t: str, lo: int = 60, hi: int = 2500) -> bool:
    n = len(t)
    return lo <= n <= hi and not _DROP_PATTERNS.search(t)


def _wsafe(name: str) -> str:
    return re.sub(r"[/\\?%*:|\"<>]", "_", name)[:60]


# ---------------------------------------------------------------- DuReader ---
def build_dureader(n_queries: int, seed: int, items: list, corpus_texts: dict) -> None:
    """dev.jsonl 抽样：3 查询的正段落 + 干扰段落 → 1 个多主题 txt 文档。"""
    rng = random.Random(seed)
    rng.shuffle(items)
    labeled_docids: set[str] = set()
    picked = []
    for it in items:
        if len(picked) >= n_queries:
            break
        q = _clean_text(it.get("query"))
        pos = it.get("positive_passages") or []
        if not (6 <= len(q) <= 40) or not pos:
            continue
        texts = [_clean_text(p.get("text")) for p in pos]
        texts = [t for t in texts if _passage_ok(t)]
        if len(texts) < 1 or sum(len(t) for t in texts) > 3000:
            continue
        picked.append((q, texts))
        for p in pos:
            labeled_docids.add(str(p.get("docid")))
    # 干扰段落：从 passage_collection 前 8 万行抽样（docid 未被任何已选查询标注）
    distractors: list[str] = []
    gz = os.path.join(RAW, "dureader_passages.tsv.gz")
    if os.path.isfile(gz):
        with gzip.open(gz, "rt", encoding="utf-8", errors="ignore") as f:
            pool = []
            for i, line in enumerate(f):
                if i >= 80000:
                    break
                docid, _, txt = line.partition("\t")
                txt = _clean_text(txt)
                if docid in labeled_docids or not _passage_ok(txt, 100, 800):
                    continue
                pool.append(txt)
                if len(pool) >= 12000:
                    break
        rng.shuffle(pool)
        distractors = pool
    # 3 个查询 → 1 个文件；每文件混入 2 段干扰
    gi = 0
    for fi in range(0, len(picked), 3):
        group = picked[fi : fi + 3]
        parts = []
        for q, texts in group:
            parts.append("\n\n".join(texts))
        for _ in range(min(2, len(distractors) - gi)):
            parts.append(distractors[gi])
            gi += 1
        fname = f"百度问答精选D{fi//3 + 1:02d}.txt"
        with io.open(os.path.join(OUT_DIR, fname), "w", encoding="utf-8") as f:
            f.write("（多主题网页问答合集）\n\n" + "\n\n".join(parts) + "\n")
        corpus_texts[fname] = "\n".join(parts)
        for q, texts in group:
            items_meta.append({
                "id": f"dur_{fi//3 + 1:02d}_{len(items_meta)}",
                "query": q,
                "kind": "positive",
                "qtype": "colloquial",
                "source": "dureader_retrieval",
                "relevant_files": [fname],
                "answer_gt": None,
                "passage_gt": texts[0][:600],
            })


# ------------------------------------------------------------------- CMRC ---
def build_cmrc(n_contexts: int, seed: int, items: list, corpus_texts: dict) -> None:
    """val 集（答案可靠）抽篇章，4 篇聚合 1 文件；每篇章取 1-2 问。"""
    import pandas as pd

    rng = random.Random(seed)
    df = pd.read_parquet(os.path.join(RAW, "cmrc_val.parquet"))
    by_ctx: dict[str, list] = {}
    for _, r in df.iterrows():
        ctx = _clean_text(r["context"])
        if not (200 <= len(ctx) <= 1500) or _DROP_PATTERNS.search(ctx):
            continue
        ans = r["answers"]
        texts = list(ans["text"]) if hasattr(ans, "keys") else []
        a = _clean_text(texts[0]) if texts else ""
        q = _clean_text(r["question"])
        if not a or not (2 <= len(a) <= 30) or len(q) < 6:
            continue
        by_ctx.setdefault(ctx, []).append((q, a))
    ctxs = [c for c, qs in by_ctx.items() if qs]
    rng.shuffle(ctxs)
    ctxs = ctxs[:n_contexts]
    for fi in range(0, len(ctxs), 4):
        group = ctxs[fi : fi + 4]
        fname = f"CMRC百科文选C{fi//4 + 1:02d}.md"
        body = "\n\n".join(f"## 节选{si + 1}\n\n{c}" for si, c in enumerate(group))
        with io.open(os.path.join(OUT_DIR, fname), "w", encoding="utf-8") as f:
            f.write("# 维基百科条目节选合集\n\n" + body + "\n")
        corpus_texts[fname] = "\n".join(group)
        for c in group:
            for q, a in by_ctx[c][:2]:
                items_meta.append({
                    "id": f"cmrc_{len(items_meta)}",
                    "query": q,
                    "kind": "positive",
                    "qtype": "fact",
                    "source": "cmrc2018",
                    "relevant_files": [fname],
                    "answer_gt": a,
                    "passage_gt": c[:600],
                })


# ------------------------------------------------------------------- 维基 ---
def build_wiki(n_articles: int, seed: int, items: list, corpus_texts: dict, use_llm: bool) -> None:
    """长条目（3000-20000 字）→ 整篇 .md；LLM 生成问答并自检答案可在原文定位。"""
    rng = random.Random(seed)
    path = os.path.join(RAW, "wiki_cn.json")
    if not os.path.isfile(path):
        print("[wiki] 原始文件缺失，跳过")
        return
    with io.open(path, encoding="utf-8") as f:
        data = json.load(f)
    arts = []
    for d in data:
        t = _clean_text(d.get("completion") or "")
        title = _clean_text(d.get("prompt") or "").strip("Question: 答案")[:30] or "条目"
        # 截到最后一个完整句
        if len(t) > 400:
            cut = t.rfind("。")
            if cut > 300:
                t = t[: cut + 1]
        if 3000 <= len(t) <= 20000 and not _DROP_PATTERNS.search(t):
            arts.append((title, t))
    rng.shuffle(arts)
    arts = arts[:n_articles]
    llm = None
    if use_llm and arts:
        from services.llm_factory import build_chat_llm
        llm = build_chat_llm(0.0)
    made = 0
    for ai, (title, t) in enumerate(arts):
        fname = f"维基长文W{ai + 1:02d}_{_wsafe(title)}.md"
        with io.open(os.path.join(OUT_DIR, fname), "w", encoding="utf-8") as f:
            f.write(f"# {title}\n\n{t}\n")
        corpus_texts[fname] = t
        if llm is None:
            continue
        prompt = (
            "以下是百科条目全文。请生成 5 个事实型问答对，要求：问题自然、答案必须是原文中"
            "连续出现的原话（10-30 字）、五个问题分布在不同段落。只输出 JSON 数组，"
            '格式 [{"q":"...","a":"..."}]，不要输出其他内容。\n\n' + t[:8000]
        )
        try:
            resp = llm.invoke(prompt)
            raw = str(getattr(resp, "content", resp))
            m = re.search(r"\[.*\]", raw, re.S)
            pairs = json.loads(m.group(0)) if m else []
        except Exception as e:  # LLM 失败不阻塞构建
            print(f"[wiki] LLM 生成失败({ai + 1}): {e}")
            pairs = []
        kept = 0
        for p in pairs:
            if kept >= 5:
                break
            q, a = _clean_text(p.get("q")), _clean_text(p.get("a"))
            if not q or not a:
                continue
            norm = re.sub(r"[\s，。、；：""''（）]", "", t)
            if re.sub(r"[\s，。、；：""''（）]", "", a) not in norm:
                continue  # 答案必须能在原文定位
            items_meta.append({
                "id": f"wiki_{ai + 1:02d}_{kept}",
                "query": q,
                "kind": "positive",
                "qtype": "fact",
                "source": "wikipedia_llm",
                "relevant_files": [fname],
                "answer_gt": a,
                "passage_gt": None,
            })
            kept += 1
        made += kept
        print(f"[wiki] {fname} 生成 {kept} 问")
    if llm:
        print(f"[wiki] 共生成 {made} 问")


# ------------------------------------------------------------ office 文件 ---
def build_office_files() -> None:
    from docx import Document
    import openpyxl
    from pptx import Presentation
    from pptx.util import Inches, Pt

    # docx：中医养生基础
    doc = Document()
    doc.add_heading("中医养生基础", 0)
    doc.add_heading("一、四季养生总则", 1)
    doc.add_paragraph(
        "中医养生的四季作息遵循《黄帝内经》\"春夏养阳、秋冬养阴\"的原则：春季晚睡早起、广步于庭，"
        "顺应生发之气；夏季夜卧早起、无厌于日；秋季早卧早起、与鸡俱兴；冬季早卧晚起、必待日光，藏阳护阴。"
    )
    doc.add_heading("二、常用代茶饮", 1)
    doc.add_paragraph(
        "枸杞菊花茶：清肝明目，适合长期用眼、易上火者；脾胃虚寒、容易腹泻者不宜多饮。"
        "山楂陈皮茶：消食化积、理气健脾，适合餐后腹胀者。酸梅汤（乌梅、山楂、甘草、桂花）：生津止渴，"
        "夏季饮用尤宜，胃酸过多者慎用。"
    )
    doc.add_heading("三、常用保健穴位", 1)
    doc.add_paragraph(
        "足三里：位于膝盖外侧凹陷（外膝眼）下四横指、胫骨外一横指处，为强壮保健要穴，"
        "常按揉可健脾和胃、增强体质，每日按压 5-10 分钟。合谷：手背第一、二掌骨之间，"
        "主头痛、牙痛。涌泉：足底前三分之一凹陷处，睡前搓揉可引火归元、助眠。"
    )
    doc.add_heading("四、饮食有节", 1)
    doc.add_paragraph(
        "《内经》谓\"五谷为养，五果为助，五畜为益，五菜为充\"。吃饭七分饱，晚餐宜清淡且不晚于睡前 3 小时；"
        "忌过食生冷（伤脾阳）与肥甘厚味（生痰湿）。体质辨识：怕冷乏力多属阳虚，宜温补；口干烦热多属阴虚，"
        "宜滋阴；形体肥胖、舌苔厚腻多属痰湿，宜清淡祛湿。"
    )
    doc.save(os.path.join(OUT_DIR, "中医养生基础_20260819.docx"))

    # xlsx：公司年度培训计划
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "2026年度培训"
    ws.append(["季度", "培训项目", "对象", "形式", "学时", "负责部门"])
    ws.append(["Q1（1-3月）", "新员工入职引导", "全体新员工", "线下集训", "24", "人力资源部"])
    ws.append(["Q1（1-3月）", "消防安全与应急演练", "全员", "线下演练", "4", "行政部"])
    ws.append(["Q2（4-6月）", "项目管理进阶（PMP实务）", "项目经理", "线上直播", "32", "项目管理办公室"])
    ws.append(["Q2（4-6月）", "高效沟通与跨部门协作", "入职1-3年员工", "线下工作坊", "12", "人力资源部"])
    ws.append(["Q3（7-8月）", "新员工入职引导", "夏季批次新员工", "线下集训", "24", "人力资源部"])
    ws.append(["Q3（7-8月）", "信息安全与数据合规", "全员", "线上必修", "8", "信息技术部"])
    ws.append(["Q4（10-12月）", "年度绩效管理与复盘方法", "各级管理者", "线下工作坊", "16", "人力资源部"])
    ws.append(["Q4（10-12月）", "人工智能工具应用实践", "全员自愿报名", "线上选修", "10", "信息技术部"])
    for col, w in zip("ABCDEF", [12, 26, 18, 12, 8, 16]):
        ws.column_dimensions[col].width = w
    wb.save(os.path.join(OUT_DIR, "公司年度培训计划_20260819.xlsx"))

    # pptx：智能手表发布
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "星环 Watch X 发布要点"
    slide.placeholders[1].text = "为健康而生的全场景智能手表"
    bullets = [
        ("核心健康功能", [
            "24 小时连续血氧监测，低血氧震动提醒",
            "心率异常（房颤）AI 筛查，通过二类医疗器械认证",
            "睡眠分期分析：浅睡/深睡/REM + 呼吸质量评分",
            "女性健康管理：周期记录与备孕体温曲线",
        ]),
        ("硬件与续航", [
            "1.43 英寸 AMOLED 屏，峰值亮度 1000 尼特",
            "典型使用续航 14 天，重度使用 7 天，磁吸快充 5 分钟用一天",
            "5ATM 防水（50 米），支持游泳与开放水域模式",
            "航空级铝合金表体 + 蓝宝石镜面，重 32 克",
        ]),
        ("价格与首发", [
            "星环 Watch X 定价 1299 元",
            "首发期下单立减 100 元，加赠氟橡胶表带一条",
            "4 月 20 日全渠道开售，前 1000 名赠一年延保",
        ]),
    ]
    for title, lines in bullets:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        body = s.placeholders[1].text_frame
        body.text = lines[0]
        for ln in lines[1:]:
            p = body.add_paragraph()
            p.text = ln
            p.font.size = Pt(18)
    prs.save(os.path.join(OUT_DIR, "智能手表星环WatchX发布_20260819.pptx"))


# ------------------------------------------------------------------ 负样本 ---
def build_negatives(corpus_texts: dict, n_pure: int) -> None:
    import eval_set_handcrafted as H

    # 纯负：查询整体与语料无关键词重叠（以查询前 6 字抽 3 个片段粗查）
    all_text = re.sub(r"\s+", "", "\n".join(corpus_texts.values()))
    kept = 0
    for q in H.NEG_PURE_CANDIDATES:
        if kept >= n_pure:
            break
        probe = re.sub(r"[\s，。、？?!！的有什么哪些怎么]", "", q)
        frags = [probe[i : i + 3] for i in range(0, max(1, len(probe) - 2), 4)][:4]
        if any(fr and fr in all_text for fr in frags):
            continue
        items_meta.append({
            "id": f"neg_pure_{kept}",
            "query": q,
            "kind": "negative",
            "subtype": "pure",
            "source": "selfbuilt",
            "relevant_files": [],
            "answer_gt": None,
            "passage_gt": None,
        })
        kept += 1
    for i, q in enumerate(H.NEG_TRAP):
        items_meta.append({
            "id": f"neg_trap_{i}",
            "query": q,
            "kind": "negative",
            "subtype": "trap",
            "source": "selfbuilt",
            "relevant_files": [],
            "answer_gt": None,
            "passage_gt": None,
        })
    return kept


# -------------------------------------------------------------------- main ---
items_meta: list = []


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--seed", type=int, default=19)
    ap.add_argument("--dureader-n", type=int, default=240)
    ap.add_argument("--cmrc-n", type=int, default=140)
    ap.add_argument("--wiki-n", type=int, default=12)
    ap.add_argument("--neg-pure-n", type=int, default=45)
    ap.add_argument("--skip-wiki-llm", action="store_true", help="跳过 LLM 问答生成（离线调试用）")
    ap.add_argument("--keep-raw", action="store_true", help="保留 3.5GB passage_collection 原始文件")
    args = ap.parse_args()

    os.makedirs(OUT_DIR, exist_ok=True)
    corpus_texts: dict[str, str] = {}
    # 已有自建文件文本（自建查询的文件引用校验用）
    for fn in os.listdir(OUT_DIR):
        if fn.endswith((".md", ".txt", ".html", ".csv")) and not fn.startswith("语料来源说明"):
            p = os.path.join(OUT_DIR, fn)
            try:
                with io.open(p, encoding="utf-8") as f:
                    corpus_texts[fn] = f.read()
            except Exception:
                pass

    print("== 生成 office 三件套 ==")
    build_office_files()

    print("== DuReader-retrieval ==")
    items = []
    with gzip.open(os.path.join(RAW, "dureader_ranking_dev.jsonl.gz"), "rt", encoding="utf-8") as f:
        for line in f:
            items.append(json.loads(line))
    build_dureader(args.dureader_n, args.seed, items, corpus_texts)

    print("== CMRC ==")
    build_cmrc(args.cmrc_n, args.seed, items, corpus_texts)

    print("== 维基长文 ==")
    build_wiki(args.wiki_n, args.seed, items, corpus_texts, use_llm=not args.skip_wiki_llm)

    print("== 自建查询（旧 104 + 新） ==")
    import eval_set_handcrafted as H

    for i, (q, files) in enumerate(H.LEGACY_EVAL_SET):
        items_meta.append({
            "id": f"legacy_{i}",
            "query": q,
            "kind": "positive" if files else "negative",
            "subtype": "legacy",
            "source": "selfbuilt_20260817",
            "relevant_files": files,
            "answer_gt": None,
            "passage_gt": None,
        })
    for i, (q, files, ans) in enumerate(H.NEW_SELFBUILT_QUERIES):
        items_meta.append({
            "id": f"selfnew_{i}",
            "query": q,
            "kind": "positive",
            "qtype": "mixed",
            "source": "selfbuilt_20260819",
            "relevant_files": files,
            "answer_gt": ans,
            "passage_gt": None,
        })

    print("== 负样本 ==")
    n_pure = build_negatives(corpus_texts, args.neg_pure_n)

    # ---- 自检 ----
    print("\n== 自检 ==")
    errs = []
    files_on_disk = set(os.listdir(OUT_DIR))
    for it in items_meta:
        for fn in it["relevant_files"]:
            if fn not in files_on_disk:
                errs.append(f"引用文件不存在: {it['id']} -> {fn}")
    queries = [it["query"] for it in items_meta]
    dup = len(queries) - len(set(queries))
    # 汇总统计
    pos = [it for it in items_meta if it["kind"] == "positive"]
    neg = [it for it in items_meta if it["kind"] == "negative"]
    by_src: dict[str, int] = {}
    for it in pos:
        by_src[it["source"]] = by_src.get(it["source"], 0) + 1
    print(f" 正样本 {len(pos)}：{by_src}")
    print(f" 负样本 {len(neg)}（纯负 {n_pure} + 陷阱 {len(H.NEG_TRAP)} + legacy {len(neg) - n_pure - len(H.NEG_TRAP)}）")
    print(f" 总条数 {len(items_meta)}，重复查询 {dup}，语料文件 {len(files_on_disk)}")
    if errs:
        print(" !! 自检失败：")
        for e in errs[:20]:
            print("   ", e)
        raise SystemExit(1)

    with io.open(OUT_JSON, "w", encoding="utf-8") as f:
        json.dump(items_meta, f, ensure_ascii=False, indent=1)
    print(f"\n写出 {OUT_JSON}（{len(items_meta)} 条）")

    # 语料说明
    with io.open(os.path.join(OUT_DIR, "语料来源说明_20260819.md"), "w", encoding="utf-8") as f:
        f.write(
            "# eval_corpus_v2 语料来源说明（2026-08-19）\n\n"
            "- 百度问答精选D*.txt：DuReader-retrieval dev 正段落拼合 + 干扰段落（ACL 2022）\n"
            "- CMRC百科文选C*.md：CMRC 2018 validation 篇章聚合（哈工大讯飞人工标注）\n"
            "- 维基长文W*.md：中文维基完整长条目（LLM 生成问答、答案原文自检）\n"
            f"- 自建 {sum(1 for x in files_on_disk if '_20260819' in x)} 篇（含 docx/xlsx/pptx）"
            " + 旧 2026-08-17 十篇迁移\n"
            f"- 评测集共 {len(items_meta)} 条（正 {len(pos)} / 负 {len(neg)}），明细见 scripts/EVAL_SET_V2.json\n"
        )

    if not args.keep_raw:
        big = os.path.join(RAW, "dureader_passages.tsv.gz")
        if os.path.isfile(big):
            os.remove(big)
            print("已删除 3.5GB passage_collection 原始文件（--keep-raw 可保留）")


if __name__ == "__main__":
    main()
