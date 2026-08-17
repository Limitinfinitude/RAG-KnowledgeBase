"""统一文档解析：扩展名 → 解析器，一处实现，多处复用。

此前同一格式的解析逻辑分散在四处的 elif 链里（file_loader 常规入库、
instant_document_loader Streamlit 即时对话、instant_doc_parse Web 即时对话、
document_preview 原文查看），加一种格式要改多处且行为逐渐发散。
现在全部收敛到这里；大文件流式入库（ingest_streaming 的分段迭代器）
因内存模型不同保持独立，但其支持的扩展名以本模块 STREAMABLE_EXTENSIONS 为准。

约定：解析器签名 (temp_path: str, original_name: str) -> List[Document]，
每个 Document 带 metadata={"source_file": 原名, "file_type": 扩展名}。
依赖在函数内延迟导入，保持本模块可被轻量引用。
"""
from __future__ import annotations

import io
import logging
import os
import subprocess
import time
from typing import Callable, Dict, List, Tuple

from langchain_core.documents import Document

logger = logging.getLogger(__name__)

# 展示顺序即默认上传白名单顺序；.doc 不在其中（见 parse_doc）
SUPPORTED_EXTENSIONS = ["pdf", "docx", "pptx", "txt", "md", "csv", "html", "xlsx", "xls",
                        "jpg", "jpeg", "png"]

# 大文件流式入库（分段迭代器）支持的扩展名；其余格式走整文件解析
STREAMABLE_EXTENSIONS = {"txt", "md", "pdf", "docx"}

SUPPORTED_EXTENSIONS_TEXT = ", ".join("." + e for e in SUPPORTED_EXTENSIONS)


# ------------------- Tesseract OCR 公共配置 -------------------

_cached_ocr_langs: List[str] | None = None


def _apply_tesseract_cmd() -> "object":
    """配置并返回 pytesseract；Tesseract 未安装时抛出带安装指引的错误。"""
    import pytesseract

    from config import TESSERACT_CMD

    if not TESSERACT_CMD:
        raise RuntimeError(
            "未找到 Tesseract OCR，扫描版 PDF / 图片无法识别。"
            "Windows: https://github.com/UB-Mannheim/tesseract/wiki 安装；"
            "Linux: apt install tesseract-ocr tesseract-ocr-chi-sim tesseract-ocr-eng；"
            "或设置环境变量 TESSERACT_CMD 指向可执行文件。"
        )
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    return pytesseract


def resolve_ocr_langs() -> str:
    """按已安装语言包选择 OCR 语言组合：中英混排优先 chi_sim+eng。

    依赖方（本机/服务器）语言包安装情况不一，写死组合会在缺包时报
    "Failed loading language"；这里查一次 --list-langs 并缓存，缺哪个降级哪个。
    """
    global _cached_ocr_langs
    if _cached_ocr_langs is None:
        try:
            pt = _apply_tesseract_cmd()
            out = subprocess.run(
                [pt.pytesseract.tesseract_cmd, "--list-langs"],
                capture_output=True, text=True, timeout=15,
            )
            langs = [ln.strip() for ln in (out.stdout or "").splitlines()[1:] if ln.strip()]
            _cached_ocr_langs = langs or ["chi_sim", "eng"]
        except Exception:  # noqa: BLE001 — 探测失败按默认组合，OCR 调用处再暴露真实错误
            _cached_ocr_langs = ["chi_sim", "eng"]
    installed = set(_cached_ocr_langs)
    preferred = [l for l in ("chi_sim", "eng") if l in installed]
    return "+".join(preferred) if preferred else _cached_ocr_langs[0]


# ------------------- 两层 OCR：本地 Tesseract → 云端（硅基流动 DeepSeek-OCR 等） -------------------

_ocr_cfg_cache: Tuple[float, Dict] | None = None
_OCR_CFG_TTL_SEC = 60.0


def _ocr_cloud_cfg() -> Dict:
    """OCR 云端配置（60s 进程内缓存，避免逐页查询设置库）。"""
    global _ocr_cfg_cache
    now = time.monotonic()
    if _ocr_cfg_cache is None or now - _ocr_cfg_cache[0] > _OCR_CFG_TTL_SEC:
        from utils.web_system_settings import get_ocr_cloud_config

        _ocr_cfg_cache = (now, get_ocr_cloud_config())
    return _ocr_cfg_cache[1]


def _build_cloud_ocr(cfg: Dict):
    from utils.siliconflow_client import SiliconFlowOCR

    return SiliconFlowOCR(api_key=cfg["api_key"], model=cfg["model"], base_url=cfg["base_url"])


def _tesseract_text_and_conf(img) -> Tuple[str, float]:
    """单次 image_to_data 同时拿文本与平均置信度（按行重组，避免跑两遍 OCR）。"""
    from pytesseract import Output

    pytesseract = _apply_tesseract_cmd()
    lang = resolve_ocr_langs()
    data = pytesseract.image_to_data(img, lang=lang, output_type=Output.DICT)
    n = len(data.get("text", []))
    confs: List[float] = []
    lines: Dict[Tuple, List[str]] = {}
    for i in range(n):
        word = str(data["text"][i] or "").strip()
        try:
            conf = float(data["conf"][i])
        except (TypeError, ValueError):
            conf = -1.0
        if conf >= 0:
            confs.append(conf)
        if not word:
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        lines.setdefault(key, []).append(word)
    text = "\n".join(" ".join(ws) for _k, ws in sorted(lines.items()))
    avg_conf = sum(confs) / len(confs) if confs else 0.0
    return text, avg_conf


def _ocr_image(img) -> Tuple[str, str]:
    """对单页/单图执行两层 OCR，返回 (文本, 引擎标记 "tesseract"|"cloud")。

    - mode=off：仅本地；
    - mode=fallback（默认）：本地先行，平均置信度低于阈值或文本过短（手写/表格/公式
      的典型表现）时该页升级走云端，云端失败保留本地结果；
    - mode=always：云端优先（限免期可用），失败回退本地。
    本地与云端全不可用时抛出原始错误（如「未找到 Tesseract」）。
    """
    cfg = _ocr_cloud_cfg()
    mode = cfg["mode"]

    tess_text, tess_conf, tess_err = "", 0.0, None
    try:
        tess_text, tess_conf = _tesseract_text_and_conf(img)
    except Exception as e:  # noqa: BLE001 — 本地失败时仍可能走云端
        tess_err = e

    def cloud_text() -> str:
        client = _build_cloud_ocr(cfg)
        with io.BytesIO() as buf:
            img.save(buf, format="PNG")
            return client.extract_text(buf.getvalue())

    if mode == "always":
        try:
            got = cloud_text()
            if got.strip():
                return got, "cloud"
        except Exception as e:  # noqa: BLE001
            logger.warning("[OCR] 云端识别失败，回退本地: %s", e)
        if tess_err is not None:
            raise tess_err
        return tess_text, "tesseract"

    if mode == "fallback":
        low_quality = (
            tess_err is not None
            or tess_conf < cfg["conf_threshold"]
            or len(tess_text.strip()) < 10
        )
        if low_quality:
            try:
                got = cloud_text()
                if got.strip():
                    logger.info("[OCR] 本地置信度 %.0f 低于阈值 %d，该页已用云端识别", tess_conf, cfg["conf_threshold"])
                    return got, "cloud"
            except Exception as e:  # noqa: BLE001
                logger.warning("[OCR] 云端回退失败，保留本地结果: %s", e)

    if tess_err is not None:
        raise tess_err
    return tess_text, "tesseract"


def detect_text_file_encoding(path: str, sample_size: int = 262144) -> str:
    """探测文本编码：utf-8-sig / utf-8 / gb18030 / gbk，兜底 utf-8。"""
    with open(path, "rb") as f:
        raw = f.read(sample_size)
    if not raw:
        return "utf-8"
    for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            raw.decode(enc, errors="strict")
            return enc
        except UnicodeDecodeError:
            continue
    return "utf-8"


def _mk_doc(text: str, original_name: str, file_type: str) -> Document:
    return Document(
        page_content=text,
        metadata={"source_file": original_name, "file_type": file_type},
    )


# ------------------- 文本类 -------------------

def parse_txt(temp_path: str, original_name: str, file_type: str = "txt") -> List[Document]:
    enc = detect_text_file_encoding(temp_path)
    with open(temp_path, "r", encoding=enc, errors="replace") as f:
        text = f.read()
    if not text.strip():
        raise ValueError("文本文件为空或未提取到内容")
    return [_mk_doc(text.strip(), original_name, file_type)]


def parse_md(temp_path: str, original_name: str) -> List[Document]:
    return parse_txt(temp_path, original_name, file_type="md")


# ------------------- Office 文档 -------------------

def parse_pdf(temp_path: str, original_name: str) -> List[Document]:
    """PDF：优先文本层提取；扫描版（无文本层）自动回退逐页 OCR。"""
    from langchain_community.document_loaders import PyPDFLoader

    docs = PyPDFLoader(temp_path).load()
    if any(d.page_content.strip() for d in docs):
        out = []
        for i, d in enumerate(docs):
            if not d.page_content.strip():
                continue
            d.metadata["source_file"] = original_name
            d.metadata.setdefault("file_type", "pdf")
            d.metadata.setdefault("page", i + 1)
            out.append(d)
        return out

    # 无文本层：渲染成图片做两层 OCR（本地 Tesseract → 云端）
    from pdf2image import convert_from_path

    logger.info("[Parser] PDF 无文本层，走 OCR：%s", original_name)
    images = convert_from_path(temp_path, fmt="png", dpi=300)
    ocr_parts: List[str] = []
    cloud_pages = 0
    for i, img in enumerate(images):
        try:
            text, engine = _ocr_image(img)
            cloud_pages += engine == "cloud"
            ocr_parts.append(f"第{i + 1}页：\n{text}")
        except Exception as e:  # noqa: BLE001 — 单页失败不应毁掉整篇
            logger.warning("第%d页OCR失败: %s", i + 1, e)
            ocr_parts.append(f"第{i + 1}页：OCR识别失败")
    if cloud_pages:
        logger.info("[Parser] PDF OCR 完成：%d 页中 %d 页由云端识别", len(images), cloud_pages)
    text = "\n\n".join(ocr_parts).strip()
    if not text:
        raise ValueError("PDF 未提取到文本（OCR 也未识别出内容）")
    return [_mk_doc(text, original_name, "pdf")]


def parse_image(temp_path: str, original_name: str, file_type: str = "jpg") -> List[Document]:
    """图片 OCR：截图 / 扫描件直接入库（jpg / jpeg / png），同样走两层 OCR。"""
    from PIL import Image

    with Image.open(temp_path) as img:
        text, _engine = _ocr_image(img)
    text = (text or "").strip()
    if not text:
        raise ValueError("图片中未识别到文字（可能是纯图形，或分辨率过低）")
    return [_mk_doc(text, original_name, file_type)]


def parse_docx(temp_path: str, original_name: str) -> List[Document]:
    from docx import Document as DocxDocument

    doc = DocxDocument(temp_path)
    parts: List[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in table.rows]
        if rows:
            parts.append("\n".join(rows))
    text = "\n\n".join(parts).strip()
    if not text:
        raise ValueError("Word 文档中未提取到文本内容")
    return [_mk_doc(text, original_name, "docx")]


def parse_doc(temp_path: str, original_name: str) -> List[Document]:
    # python-docx / docx2txt 均不支持老版二进制 .doc；给出可操作的提示而不是误导性报错
    raise ValueError(
        "不支持老版 .doc 格式（Word 97-2003 二进制）。"
        "请用 Word/WPS 打开后「另存为」.docx 再上传。"
    )


def parse_pptx(temp_path: str, original_name: str) -> List[Document]:
    from pptx import Presentation

    prs = Presentation(temp_path)
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: List[str] = [f"幻灯片 {i}："]
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    lines.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        except Exception:  # noqa: BLE001
            notes = ""
        if notes:
            lines.append(f"[备注] {notes}")
        if len(lines) > 1:
            parts.append("\n".join(lines))
    text = "\n\n".join(parts).strip()
    if not text:
        raise ValueError("PPT 中未提取到文本内容")
    return [_mk_doc(text, original_name, "pptx")]


# ------------------- 表格类 -------------------

def _df_to_text(df) -> str:
    import pandas as pd

    return df.fillna("").to_string(index=False)


def parse_xlsx(temp_path: str, original_name: str, file_type: str = "xlsx") -> List[Document]:
    import pandas as pd

    excel_file = pd.ExcelFile(temp_path)
    sheets = []
    for sheet_name in excel_file.sheet_names:
        sheet_text = f"工作表: {sheet_name}\n\n"
        sheet_text += _df_to_text(pd.read_excel(excel_file, sheet_name=sheet_name))
        sheets.append(sheet_text)
    text = "\n\n" + "=" * 50 + "\n\n".join(sheets)
    if not text.strip():
        raise ValueError("Excel 中未提取到文本内容")
    return [_mk_doc(text, original_name, file_type)]


def parse_xls(temp_path: str, original_name: str) -> List[Document]:
    return parse_xlsx(temp_path, original_name, file_type="xls")


def parse_csv(temp_path: str, original_name: str) -> List[Document]:
    import pandas as pd

    enc = detect_text_file_encoding(temp_path)
    if enc == "utf-8":  # utf-8-sig 探测通过但无 BOM 时返回 utf-8；统一用 utf-8-sig 兼容 BOM
        enc = "utf-8-sig"
    df = pd.read_csv(temp_path, encoding=enc)
    text = _df_to_text(df)
    if not text.strip():
        raise ValueError("CSV 中未提取到内容")
    return [_mk_doc(text, original_name, "csv")]


# ------------------- 网页类 -------------------

def parse_html(temp_path: str, original_name: str) -> List[Document]:
    from bs4 import BeautifulSoup

    with open(temp_path, "rb") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = (soup.title.string or "").strip() if soup.title and soup.title.string else ""
    text = soup.get_text("\n")
    # 压缩 3 连以上空行
    lines = [ln.strip() for ln in text.splitlines()]
    packed: List[str] = []
    blank = 0
    for ln in lines:
        if ln:
            packed.append(ln)
            blank = 0
        elif blank < 1:
            packed.append("")
            blank += 1
    body = "\n".join(packed).strip()
    if not body:
        raise ValueError("HTML 中未提取到文本内容")
    if title:
        body = f"标题: {title}\n\n{body}"
    return [_mk_doc(body, original_name, "html")]


# ------------------- 注册表与入口 -------------------

ParserFn = Callable[[str, str], List[Document]]

PARSERS: Dict[str, ParserFn] = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "doc": parse_doc,  # 在注册表中以给出定向报错，但不在 SUPPORTED_EXTENSIONS 白名单内
    "pptx": parse_pptx,
    "txt": parse_txt,
    "md": parse_md,
    "csv": parse_csv,
    "html": parse_html,
    "xlsx": parse_xlsx,
    "xls": parse_xls,
    "jpg": parse_image,
    "jpeg": lambda p, n: parse_image(p, n, file_type="jpeg"),
    "png": lambda p, n: parse_image(p, n, file_type="png"),
}


def parse_file_to_documents(temp_path: str, original_name: str) -> List[Document]:
    """按扩展名解析整个文件为 Document 列表（常规入库 / 即时对话 / 原文查看共用）。"""
    raw = original_name or temp_path
    ext = os.path.splitext(raw.replace("\\", "/"))[1].lower().lstrip(".")
    parser = PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"不支持的文件类型: .{ext}，支持: {SUPPORTED_EXTENSIONS_TEXT}")
    return parser(temp_path, original_name or os.path.basename(temp_path))
