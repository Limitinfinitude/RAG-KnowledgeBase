"""统一文档解析器（utils/document_parsers）单测：新格式解析 + 白名单一致性 + .doc 定向报错。"""
from __future__ import annotations

import pytest

from utils.document_parsers import SUPPORTED_EXTENSIONS, parse_file_to_documents


def _write(tmp_path, name: str, data: bytes) -> str:
    p = tmp_path / name
    p.write_bytes(data)
    return str(p)


def test_txt_utf8_and_gb18030(tmp_path):
    p8 = _write(tmp_path, "a.txt", "你好世界".encode("utf-8"))
    docs = parse_file_to_documents(p8, "a.txt")
    assert "你好世界" in docs[0].page_content
    assert docs[0].metadata["file_type"] == "txt"
    assert docs[0].metadata["source_file"] == "a.txt"

    pg = _write(tmp_path, "b.txt", "中文编码测试".encode("gb18030"))
    docs = parse_file_to_documents(pg, "b.txt")
    assert "中文编码测试" in docs[0].page_content


def test_md(tmp_path):
    p = _write(tmp_path, "note.md", "# 标题\n\n正文".encode("utf-8"))
    docs = parse_file_to_documents(p, "note.md")
    assert docs[0].metadata["file_type"] == "md"


def test_csv_chinese_with_bom(tmp_path):
    p = _write(tmp_path, "t.csv", "名称,数量\n苹果,3\n香蕉,5\n".encode("utf-8-sig"))
    docs = parse_file_to_documents(p, "t.csv")
    text = docs[0].page_content
    assert "苹果" in text and "香蕉" in text
    assert docs[0].metadata["file_type"] == "csv"


def test_html_strips_script_style_keeps_title(tmp_path):
    html = (
        "<html><head><title>测试页</title><style>.x{color:red}</style>"
        "<script>var a=1;</script></head>"
        "<body><h1>主标题</h1><p>正文内容</p></body></html>"
    ).encode("utf-8")
    p = _write(tmp_path, "page.html", html)
    docs = parse_file_to_documents(p, "page.html")
    text = docs[0].page_content
    assert "主标题" in text and "正文内容" in text
    assert "var a=1" not in text and "color:red" not in text
    assert "测试页" in text
    assert docs[0].metadata["file_type"] == "html"


def test_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "检索增强生成"
    slide.placeholders[1].text = "RAG 的三个阶段"
    prs.save(str(tmp_path / "deck.pptx"))

    docs = parse_file_to_documents(str(tmp_path / "deck.pptx"), "deck.pptx")
    text = docs[0].page_content
    assert "检索增强生成" in text and "RAG 的三个阶段" in text
    assert "幻灯片 1" in text
    assert docs[0].metadata["file_type"] == "pptx"


def test_docx_with_table(tmp_path):
    from docx import Document as Docx

    d = Docx()
    d.add_paragraph("第一段内容")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "键"
    t.cell(0, 1).text = "值"
    d.save(str(tmp_path / "doc.docx"))

    docs = parse_file_to_documents(str(tmp_path / "doc.docx"), "doc.docx")
    text = docs[0].page_content
    assert "第一段内容" in text
    assert "键 | 值" in text


def test_xlsx(tmp_path):
    import pandas as pd

    p = str(tmp_path / "s.xlsx")
    pd.DataFrame({"列A": [1, 2], "列B": ["x", "y"]}).to_excel(p, index=False)
    docs = parse_file_to_documents(p, "s.xlsx")
    assert "列A" in docs[0].page_content
    assert "工作表" in docs[0].page_content


def test_doc_rejected_with_actionable_hint(tmp_path):
    p = _write(tmp_path, "old.doc", b"\xd0\xcf\x11\xe0")  # OLE2 头（真实 .doc 特征）
    with pytest.raises(ValueError, match="另存为"):
        parse_file_to_documents(p, "old.doc")


def test_unsupported_ext_lists_supported(tmp_path):
    p = _write(tmp_path, "x.foo", b"data")
    with pytest.raises(ValueError, match="不支持的文件类型"):
        parse_file_to_documents(p, "x.foo")


def test_whitelist_consistency():
    from utils.web_system_settings import _DEFAULT_EXT
    from utils.instant_doc_parse import INSTANT_ALLOWED_EXT

    assert _DEFAULT_EXT == SUPPORTED_EXTENSIONS
    assert INSTANT_ALLOWED_EXT == {"." + e for e in SUPPORTED_EXTENSIONS}


def test_instant_parse_bytes_end_to_end():
    from utils.instant_doc_parse import parse_upload_bytes

    text = parse_upload_bytes("t.csv", "名称,数量\n苹果,3\n".encode("utf-8"))
    assert "苹果" in text

    with pytest.raises(ValueError, match="另存为"):
        parse_upload_bytes("old.doc", b"\xd0\xcf\x11\xe0")

    with pytest.raises(ValueError, match="不支持的格式"):
        parse_upload_bytes("x.foo", b"data")


# ------------------- OCR：语言探测与图片解析 -------------------

def _reset_lang_cache():
    import utils.document_parsers as dp
    dp._cached_ocr_langs = None


def test_resolve_ocr_langs_prefers_chi_sim_eng(monkeypatch):
    import utils.document_parsers as dp

    class FakeProc:
        def __init__(self, stdout):
            self.stdout = stdout

    monkeypatch.setattr(dp.subprocess, "run", lambda *a, **k: FakeProc("List of available languages:\neng\nchi_sim\nosd\n"))
    _reset_lang_cache()
    assert dp.resolve_ocr_langs() == "chi_sim+eng"

    monkeypatch.setattr(dp.subprocess, "run", lambda *a, **k: FakeProc("List of available languages:\neng\nosd\n"))
    _reset_lang_cache()
    assert dp.resolve_ocr_langs() == "eng"

    monkeypatch.setattr(dp.subprocess, "run", lambda *a, **k: FakeProc("List of available languages:\nchi_sim\n"))
    _reset_lang_cache()
    assert dp.resolve_ocr_langs() == "chi_sim"

    # 探测失败（如 tesseract 报错）按默认组合，不抛异常
    def boom(*a, **k):
        raise OSError("no tesseract")
    monkeypatch.setattr(dp.subprocess, "run", boom)
    _reset_lang_cache()
    assert dp.resolve_ocr_langs() == "chi_sim+eng"
    _reset_lang_cache()


def _tesseract_available() -> bool:
    from config import TESSERACT_CMD
    return bool(TESSERACT_CMD)


@pytest.mark.skipif(not _tesseract_available(), reason="本机无 Tesseract，跳过真图 OCR")
def test_parse_image_real_ocr(tmp_path):
    from PIL import Image, ImageDraw, ImageFont

    font = None
    for name in ("arial.ttf", "DejaVuSans.ttf", "C:/Windows/Fonts/arial.ttf"):
        try:
            font = ImageFont.truetype(name, 48)
            break
        except OSError:
            continue
    if font is None:
        pytest.skip("无可用 TTF 字体")

    img = Image.new("RGB", (640, 160), "white")
    d = ImageDraw.Draw(img)
    d.text((30, 50), "HELLO RAG 2468", fill="black", font=font)
    p = tmp_path / "shot.png"
    img.save(str(p))

    docs = parse_file_to_documents(str(p), "shot.png")
    text = docs[0].page_content.upper()
    assert "RAG" in text and "2468" in text, f"OCR 结果异常: {text!r}"
    assert docs[0].metadata["file_type"] == "png"


@pytest.mark.skipif(not _tesseract_available(), reason="本机无 Tesseract，跳过真图 OCR")
def test_parse_image_blank_raises(tmp_path, monkeypatch):
    from PIL import Image

    # 空白图会触发云端回退（且云端可能对空白图也返回内容），锁定 off 只测本地路径
    monkeypatch.setattr(dp, "_ocr_cloud_cfg", lambda: {"mode": "off", "conf_threshold": 60})
    p = tmp_path / "blank.jpg"
    Image.new("RGB", (100, 100), "white").save(str(p))
    with pytest.raises(ValueError, match="未识别到文字"):
        parse_file_to_documents(str(p), "blank.jpg")


# ------------------- 两层 OCR 回退策略（全 mock，不打真实 API） -------------------

import utils.document_parsers as dp


class _FakeCloud:
    def __init__(self, result=None, error=None):
        self.result, self.error = result, error
        self.calls = 0

    def extract_text(self, image_bytes, mime="image/png"):
        self.calls += 1
        if self.error:
            raise self.error
        return self.result


def _setup_policy(monkeypatch, *, mode, conf=90.0, text="清晰文本结果", threshold=60, cloud=None, tess_err=None):
    fake_cfg = {"mode": mode, "model": "deepseek-ai/DeepSeek-OCR", "api_key": "sk-test",
                "base_url": "https://api.siliconflow.cn", "conf_threshold": threshold}
    monkeypatch.setattr(dp, "_ocr_cloud_cfg", lambda: fake_cfg)
    monkeypatch.setattr(dp, "_tesseract_text_and_conf", lambda img: (_raise_or(tess_err), conf) if tess_err else (text, conf))
    fake_client = cloud or _FakeCloud(result="云端识别结果")
    monkeypatch.setattr(dp, "_build_cloud_ocr", lambda cfg: fake_client)
    return fake_client


def _raise_or(err):
    raise err


class _FakeImg:
    def save(self, buf, format="PNG"):
        buf.write(b"png")


def test_ocr_policy_off_never_calls_cloud(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="off", conf=10.0, text="糊")
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "tesseract" and text == "糊"
    assert cloud.calls == 0


def test_ocr_policy_fallback_low_conf_uses_cloud(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="fallback", conf=35.0, text="模糊结果")
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "cloud" and text == "云端识别结果"
    assert cloud.calls == 1


def test_ocr_policy_fallback_high_conf_skips_cloud(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="fallback", conf=95.0, text="本地高质量识别结果文本内容")
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "tesseract" and cloud.calls == 0


def test_ocr_policy_fallback_short_text_uses_cloud(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="fallback", conf=95.0, text="短")
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "cloud" and cloud.calls == 1


def test_ocr_policy_fallback_cloud_error_keeps_local(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="fallback", conf=20.0, text="本地凑合结果",
                          cloud=_FakeCloud(error=RuntimeError("api down")))
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "tesseract" and text == "本地凑合结果"
    assert cloud.calls == 1


def test_ocr_policy_fallback_tesseract_missing_cloud_saves(monkeypatch):
    """服务器没装 Tesseract 时，云端应接管而不是失败。"""
    cloud = _setup_policy(monkeypatch, mode="fallback", tess_err=RuntimeError("未找到 Tesseract OCR"))
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "cloud" and text == "云端识别结果"


def test_ocr_policy_always_cloud_first(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="always", conf=95.0, text="本地不会用到")
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "cloud" and cloud.calls == 1


def test_ocr_policy_always_cloud_error_falls_back_local(monkeypatch):
    cloud = _setup_policy(monkeypatch, mode="always", conf=95.0, text="本地兜底文本",
                          cloud=_FakeCloud(error=RuntimeError("timeout")))
    text, engine = dp._ocr_image(_FakeImg())
    assert engine == "tesseract" and text == "本地兜底文本"
